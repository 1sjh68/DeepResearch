# utils/file_handler.py

import csv
import io
import logging
import os
from collections.abc import Callable, Iterator
from typing import Any

from config.config import Config
from utils.progress import create_progress_bar

# 文件类型检测（基于内容而非扩展名）
try:
    import magic

    MAGIC_AVAILABLE = True
except ImportError:
    MAGIC_AVAILABLE = False
    logging.warning("python-magic not installed. File type detection will use extension only.")

# Define the set of supported file extensions for automatic directory scanning
SUPPORTED_EXTENSIONS: set[str] = {
    ".txt",
    ".md",
    ".py",
    ".js",
    ".html",
    ".css",
    ".json",
    ".xml",
    ".csv",
    ".pdf",
    ".docx",
    ".pptx",
    ".ppt",
}

TESSERACT_AVAILABLE = False
pytesseract: Any | None = None
convert_from_path: Any | None = None
try:
    import pytesseract as _pytesseract
    from pdf2image import convert_from_path as _convert_from_path

    pytesseract = _pytesseract
    convert_from_path = _convert_from_path
    TESSERACT_AVAILABLE = True
    logging.info("Tesseract-OCR and pdf2image found. OCR functionality is enabled.")
except ImportError:
    logging.warning("Tesseract-OCR or pdf2image not found. OCR capabilities will be disabled.")


def _read_txt(file_path: str) -> str:
    """从.txt文件读取内容。"""
    try:
        with open(file_path, encoding="utf-8", errors="ignore") as f:
            return f.read()
    except FileNotFoundError:
        raise FileNotFoundError(f"❌ 文件不存在: {file_path}\n💡 建议: 请检查文件路径是否正确")
    except PermissionError:
        raise PermissionError(f"❌ 没有权限读取文件: {file_path}\n💡 建议: 请检查文件权限或关闭占用该文件的程序")
    except Exception as e:
        raise RuntimeError(f"❌ 读取文件失败: {file_path}\n原因: {type(e).__name__}: {str(e)}\n💡 建议: 检查文件是否损坏或编码是否正确")


def _read_docx(file_path: str) -> str:
    """从.docx文件读取内容。"""
    try:
        import docx
    except ImportError:
        error_msg = f"❌ 缺少依赖: python-docx 未安装\n文件: {file_path}\n💡 解决方案: 请运行 'pip install python-docx'"
        logging.error(error_msg)
        return f"[Error: python-docx not installed to read {os.path.basename(file_path)}]"

    try:
        doc = docx.Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs if para.text])
    except Exception as e:
        error_msg = f"❌ 读取 DOCX 文件失败: {os.path.basename(file_path)}\n原因: {type(e).__name__}: {str(e)}\n💡 建议: \n  1. 检查文件是否损坏\n  2. 确认是否为有效的 .docx 格式（不是 .doc）\n  3. 尝试用 Word 打开并重新保存"
        logging.error(error_msg)
        raise RuntimeError(error_msg)


def _read_pptx(file_path: str) -> str:
    """从.pptx文件的所有幻灯片读取文本内容。"""
    try:
        import pptx
    except ImportError:
        error_msg = f"❌ 缺少依赖: python-pptx 未安装\n文件: {file_path}\n💡 解决方案: 请运行 'pip install python-pptx'"
        logging.error(error_msg)
        return f"[Error: python-pptx not installed to read {os.path.basename(file_path)}]"

    try:
        prs = pptx.Presentation(file_path)
        text_runs: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                text_frame = getattr(shape, "text_frame", None)
                if text_frame is None:
                    continue
                for paragraph in getattr(text_frame, "paragraphs", []):
                    for run in getattr(paragraph, "runs", []):
                        text_runs.append(run.text)
        return "\n".join(text_runs)
    except Exception as e:
        error_msg = f"❌ 读取 PPTX 文件失败: {os.path.basename(file_path)}\n原因: {type(e).__name__}: {str(e)}\n💡 建议: \n  1. 检查文件是否损坏\n  2. 确认是否为有效的 .pptx 格式（不是 .ppt）\n  3. 尝试用 PowerPoint 打开并重新保存"
        logging.error(error_msg)
        raise RuntimeError(error_msg)


def detect_file_type_by_content(file_path: str) -> str:
    """基于文件内容检测文件类型，防止扩展名伪装"""
    if not MAGIC_AVAILABLE:
        # 回退到扩展名检测
        return os.path.splitext(file_path)[1].lower()

    try:
        mime = magic.Magic(mime=True)  # type: ignore[possibly-unbound]
        mime_type = mime.from_file(file_path)

        # 将MIME类型映射到扩展名
        mime_to_ext = {
            "text/plain": ".txt",
            "text/markdown": ".md",
            "text/x-python": ".py",
            "text/javascript": ".js",
            "application/javascript": ".js",
            "text/html": ".html",
            "text/css": ".css",
            "application/json": ".json",
            "application/xml": ".xml",
            "text/xml": ".xml",
            "text/csv": ".csv",
            "application/pdf": ".pdf",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
            "application/vnd.ms-powerpoint": ".ppt",
        }

        return mime_to_ext.get(mime_type, os.path.splitext(file_path)[1].lower())
    except Exception as e:
        logging.debug(f"Content-based type detection failed for {file_path}: {e}")
        return os.path.splitext(file_path)[1].lower()


def _read_pdf_with_ocr_fallback(file_path: str) -> str:
    """
    Intelligently reads a .pdf file, falling back to OCR if direct text extraction yields little content.
    """
    try:
        import fitz  # PyMuPDF
    except ImportError:
        error_msg = f"❌ 缺少依赖: PyMuPDF 未安装\n文件: {file_path}\n💡 解决方案: 请运行 'pip install PyMuPDF'"
        logging.error(error_msg)
        return f"[Error: PyMuPDF not installed to read {os.path.basename(file_path)}]"

    text = ""
    try:
        with fitz.open(file_path) as doc:
            page_texts: list[str] = []
            for page in doc:
                get_text_fn = getattr(page, "get_text", None)
                if callable(get_text_fn):
                    page_text = str(get_text_fn())
                    page_texts.append(page_text)
            text = "".join(page_texts).strip()
            page_count = len(doc)

        is_scanned = len(text) < 50 * page_count
        if is_scanned:
            if TESSERACT_AVAILABLE and pytesseract and convert_from_path:
                logging.info(f"  - Attempting OCR for {os.path.basename(file_path)}...")
                try:
                    # 限制处理页数，避免内存溢出
                    MAX_OCR_PAGES = 50
                    images = convert_from_path(file_path)
                    if len(images) > MAX_OCR_PAGES:
                        logging.warning(f"⚠️  PDF页数过多 ({len(images)}页)，仅处理前{MAX_OCR_PAGES}页\n💡 建议: 如需处理完整文档，请考虑拆分PDF或增加MAX_OCR_PAGES限制")
                        images = images[:MAX_OCR_PAGES]
                    # noqa: W293
                    ocr_texts = [pytesseract.image_to_string(image, lang="chi_sim+eng") for image in images]
                    text = "\n".join(ocr_texts)
                    logging.info(f"  - OCR successful for {os.path.basename(file_path)} ({len(images)}页).")
                except Exception as ocr_error:
                    logging.error(f"  - OCR failed for {os.path.basename(file_path)}: {ocr_error}")
                    text += "\n\n[OCR FAILED]"
            else:
                logging.warning(f"  - Scanned PDF detected, but Tesseract/pdftools are not available. Skipping OCR for {os.path.basename(file_path)}.")
                text += "\n\n[OCR SKIPPED]"

    except Exception as e:
        error_msg = f"❌ 处理 PDF 文件失败: {os.path.basename(file_path)}\n原因: {type(e).__name__}: {str(e)}\n💡 建议: \n  1. 检查 PDF 文件是否损坏\n  2. 确认 PDF 没有密码保护\n  3. 如果是扫描版 PDF，请安装 Tesseract-OCR\n  4. 尝试用 PDF 阅读器打开并重新保存"
        logging.error(error_msg, exc_info=True)
        return f"[Error reading PDF: {os.path.basename(file_path)}]"
    return text


def load_external_data(config: Config, file_paths: list[str]) -> str:
    """
    Loads text content from a list of file paths, supporting various formats.
    """
    if not file_paths:
        return ""

    readers: dict[str, Callable[[str], str]] = {
        ".txt": _read_txt,
        ".md": _read_txt,
        ".py": _read_txt,
        ".js": _read_txt,
        ".html": _read_txt,
        ".css": _read_txt,
        ".json": _read_txt,
        ".xml": _read_txt,
        ".csv": _read_txt,
        ".pdf": _read_pdf_with_ocr_fallback,
        ".docx": _read_docx,
        ".pptx": _read_pptx,
        ".ppt": _read_pptx,
    }

    resolved_files: list[str] = _collect_supported_paths(file_paths)

    all_content: list[str] = []
    # 使用进度条显示文件处理进度
    for fp in create_progress_bar(resolved_files, desc="📂 加载文件", unit="个文件"):
        ext = os.path.splitext(fp)[1].lower()
        if ext in readers:
            try:
                logging.info(f"Reading {ext.upper()} file: {fp}...")
                content = readers[ext](fp)
                file_header = f"\n\n--- Start of file: {os.path.basename(fp)} ---\n"
                file_footer = f"\n--- End of file: {os.path.basename(fp)} ---\n\n"
                all_content.append(file_header + content + file_footer)
            except Exception as e:
                error_msg = f"❌ 读取文件失败: {os.path.basename(fp)}\n文件路径: {fp}\n文件类型: {ext.upper()}\n错误类型: {type(e).__name__}\n错误详情: {str(e)}\n💡 建议: \n  1. 检查文件是否存在且没有损坏\n  2. 确认有读取权限\n  3. 如果是特殊格式，请安装对应的依赖库"
                logging.error(error_msg)
        else:
            logging.warning(f"Unsupported file type: {ext} for file {fp}. Skipped.")

    return "\n".join(all_content)


def _collect_supported_paths(input_paths: list[str]) -> list[str]:
    """
    Given a list of filesystem paths, return an ordered list of readable file
    paths. Directories are traversed recursively for files with supported
    extensions. Duplicate files are removed while preserving order.
    """
    collected: list[str] = []
    seen: set[str] = set()

    for raw_path in input_paths:
        if not raw_path:
            continue
        normalized_path = raw_path.replace("\\", "/")

        if not os.path.exists(normalized_path):
            logging.warning(f"⚠️  路径不存在，已跳过: '{normalized_path}'\n💡 建议: 检查路径是否拼写正确，或者文件是否已被移动/删除")
            continue

        if os.path.isdir(normalized_path):
            logging.info(f"Directory detected. Scanning for readable files in: '{normalized_path}'")
            dir_files: list[str] = []
            for root, _, files in os.walk(normalized_path):
                for file in files:
                    ext = os.path.splitext(file)[1].lower()
                    if ext in SUPPORTED_EXTENSIONS:
                        full_path = os.path.join(root, file).replace("\\", "/")
                        if full_path not in seen:
                            dir_files.append(full_path)
                            seen.add(full_path)
            logging.info("  - Found %s supported files in directory '%s'.", len(dir_files), normalized_path)
            collected.extend(dir_files)
            continue

        # 使用基于内容的类型检测
        detected_ext = detect_file_type_by_content(normalized_path)
        if detected_ext in SUPPORTED_EXTENSIONS:
            if normalized_path not in seen:
                collected.append(normalized_path)
                seen.add(normalized_path)
        else:
            logging.warning(f"⚠️  不支持的文件类型，已跳过: '{normalized_path}'\n检测到的类型: {detected_ext}\n💡 支持的格式: {', '.join(sorted(SUPPORTED_EXTENSIONS))}")

    return collected


def parse_and_validate_paths(path_string: str) -> list[str]:
    """
    Parses a comma-separated string of paths, validates them, and recursively
    collects readable files from any directories.
    """
    if not path_string or not path_string.strip():
        return []

    corrected_path_string = path_string.replace("\\", "/").strip().strip('"')

    try:
        string_reader = io.StringIO(corrected_path_string)
        path_reader: Iterator[list[str]] = csv.reader(
            string_reader,
            delimiter=",",
            quotechar='"',
            skipinitialspace=True,
        )
        potential_paths = next(path_reader)
    except (StopIteration, csv.Error):
        potential_paths = [p.strip() for p in corrected_path_string.split(",") if p.strip()]

    normalized_candidates: list[str] = [path.strip().strip('"') for path in potential_paths if path.strip()]
    return _collect_supported_paths(normalized_candidates)
